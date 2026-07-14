#!/usr/bin/env python3
"""Generate one-page PPT for sqlBotAi platform summary."""

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# Colors
PRIMARY = RGBColor(0x1A, 0x56, 0x9C)      # deep blue
ACCENT = RGBColor(0x00, 0x96, 0xD6)       # bright blue
LIGHT_BG = RGBColor(0xF0, 0xF6, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)


def set_run_font(run, size=11, bold=False, color=DARK, name="微软雅黑"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def set_paragraph_text(paragraph, text, size=11, bold=False, color=DARK):
    paragraph.text = text
    for run in paragraph.runs:
        set_run_font(run, size=size, bold=bold, color=color)


def add_header(slide, prs):
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(9.5), Inches(0.55))
    tf = title_box.text_frame
    set_paragraph_text(tf.paragraphs[0], "AI 数据资产平台（sqlBotAi）建设一页方案", size=26, bold=True, color=WHITE)

    sub_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(9.5), Inches(0.3))
    tf2 = sub_box.text_frame
    set_paragraph_text(
        tf2.paragraphs[0],
        f"问口径能答 · 问数据能查  |  {date.today().strftime('%Y.%m.%d')}",
        size=11,
        color=RGBColor(0xCC, 0xE5, 0xFF),
    )


def add_section_card(slide, left, top, width, height, title, bullets, accent=ACCENT):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xD0, 0xE4, 0xF5)
    card.line.width = Pt(1)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.08), width - Inches(0.25), Inches(0.35))
    set_paragraph_text(title_box.text_frame.paragraphs[0], title, size=13, bold=True, color=PRIMARY)

    body_box = slide.shapes.add_textbox(
        left + Inches(0.18), top + Inches(0.42), width - Inches(0.25), height - Inches(0.5)
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        set_paragraph_text(p, f"• {bullet}", size=9.5, color=DARK)
        p.space_after = Pt(3)


def add_footer(slide, prs):
    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(9.2), Inches(0.28))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_text(
        p,
        "技术栈：Spring Boot 3.2 · LangChain4j · DeepSeek · DashScope · H2 · MySQL ADS · Python 查数沙箱 · Nginx/systemd",
        size=8.5,
        color=GRAY,
    )


def build_ppt() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    add_header(slide, prs)

    # Row 1: 分析思路 + 分析过程（重点突出）
    add_section_card(
        slide,
        Inches(0.35), Inches(1.15), Inches(4.55), Inches(2.75),
        "① 搭建分析思路",
        [
            "痛点定位：数据资产知识碎片化（指标、表、血缘、SQL、经验），业务自助查数门槛高、口径偏差风险大",
            "建设愿景：打造「知识问答 + 受控查数」一体化平台，实现「问口径即答、问数据即得」",
            "核心思路：知识治理先行（指标资产 + 领域参考）→ 受控 SQL 生成（避免幻觉）→ 混合 RAG 检索 → 沙箱安全执行",
            "MVP 策略：线下营销域「进攻/轮动专案」试点，快速打通端到端闭环，再向多业务域扩展",
        ],
    )
    add_section_card(
        slide,
        Inches(5.1), Inches(1.15), Inches(4.55), Inches(2.75),
        "② 分析过程（5 阶段闭环）",
        [
            "阶段一：Wiki 知识盘点与结构化（176+ 文档，提取指标/表/血缘/领域参考）",
            "阶段二：Skill 提示词工程（意图识别 + mvp-questions + mysql-dialect 方言映射）",
            "阶段三：混合检索增强（DashScope 向量 + 关键词 + 业务同义词 + 领域强制注入）",
            "阶段四：试点场景验证（大区奖励、省区达成率、Top10、区域奖励 4 类问题全流程打通）",
            "阶段五：部署与图谱落地（阿里云 ECS 一键打包 + 知识图谱可视化 + 生产就绪）",
        ],
        accent=RGBColor(0x2E, 0x86, 0xAB),
    )

    # Row 2: 技术栈 + 落地方案 + 困难点
    col_w = Inches(3.0)
    row2_top = Inches(4.05)
    row2_h = Inches(2.85)

    add_section_card(
        slide,
        Inches(0.35), row2_top, col_w, row2_h,
        "③ 技术栈选型",
        [
            "后端：Spring Boot 3.2 + Java 17，快速交付企业级 Web 应用",
            "前端：Thymeleaf + Bootstrap，轻量后台无需独立 SPA",
            "智能：LangChain4j 统一接入 DeepSeek（问答/SQL/总结）",
            "检索：DashScope text-embedding-v2，TF-IDF 回退",
            "数据：H2 存元数据，MySQL ADS 查业务明细",
            "执行：Python query_mysql.py 只读沙箱；Nginx + systemd 部署",
        ],
        accent=RGBColor(0x27, 0xAE, 0x60),
    )
    add_section_card(
        slide,
        Inches(3.5), row2_top, col_w, row2_h,
        "④ 最终落地方案",
        [
            "统一入口 /knowledge-base：四通路路由（查数→RAG→Wiki→回退）",
            "Wiki 浏览 + 知识图谱可视化，支撑资产发现与关系理解",
            "业务查数：自然语言 → 受控 SQL → 表格结果 + 中文分析 + 来源溯源",
            "线下营销域试点上线，覆盖最终奖励、达成率、区域分析",
            "生产部署：/opt/sqlbotai，Nginx:80 反代，systemd 守护，本地 MySQL 连接",
        ],
        accent=ACCENT,
    )
    add_section_card(
        slide,
        Inches(6.65), row2_top, col_w, row2_h,
        "⑤ 困难点与应对",
        [
            "SQL 准确性：字段臆造、同义词未对齐 → Wiki 强制注入 + Skill 规则 + 方言对照",
            "口径翻译：Hive/MaxCompute → MySQL 8.4 → mysql-dialect 映射与 schema 校验",
            "粒度关联：奖励表与区域表粒度不同 → 先收敛再 JOIN，防行膨胀",
            "上下文有限：多文档竞争 4000 字上限 → 领域强制路由 + 同义词扩展检索",
            "安全边界：无鉴权、LLM 生成 SQL → 只读沙箱 + 行数/超时限制（待补 SSO）",
            "图表未贯通：render_chart.py 已备，Java 流程待接入",
        ],
        accent=ORANGE,
    )

    add_footer(slide, prs)
    return prs


def main():
    output = Path(__file__).resolve().parent / "sqlBotAi-平台建设一页方案.pptx"
    prs = build_ppt()
    prs.save(output)
    print(f"已生成: {output}")


if __name__ == "__main__":
    main()
